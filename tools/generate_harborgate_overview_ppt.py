from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "HarborGate-to-HarborBeacon-overview.pptx"

FONT = "Microsoft YaHei"
FONT_ALT = "DengXian"

BG = "F6F1E8"
CARD = "FFFDFC"
TEXT = "1F2933"
MUTED = "5D6B78"
TEAL = "0F766E"
TEAL_SOFT = "D9EFEA"
RUST = "B85C38"
RUST_SOFT = "F3E0D6"
GOLD = "BF8B30"
GOLD_SOFT = "F6E7C4"
LINE = "D7CEC0"
INK = "213547"
NAVY = "284B63"
NAVY_SOFT = "DCE7EF"
GREEN_SOFT = "DCEBDD"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 15


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_style(run, size: int, *, color: str = TEXT, bold: bool = False, name: str = FONT):
    font = run.font
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), name)
    rpr.set(qn("a:latin"), name)


def add_background(slide, color: str = BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, title: str, subtitle: str | None = None, *, page: int):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.42), Inches(1.25), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(TEAL)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.58), Inches(8.9), Inches(0.78))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run_style(run, 25, color=INK, bold=True)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = subtitle
        set_run_style(run2, 11, color=MUTED, name=FONT_ALT)

    foot = slide.shapes.add_textbox(Inches(11.55), Inches(7.0), Inches(1.0), Inches(0.22))
    tf2 = foot.text_frame
    tf2.clear()
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    run3 = p3.add_run()
    run3.text = f"{page:02d}/{TOTAL_SLIDES:02d}"
    set_run_style(run3, 9, color=MUTED, name=FONT_ALT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(6.9), Inches(11.9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(LINE)
    line.line.fill.background()


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 18,
    color: str = TEXT,
    bold_first: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_run_style(run, size, color=color, bold=bold_first and idx == 0)
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: list[str],
    *,
    accent: str = TEAL,
    title_size: int = 16,
    body_size: int = 13,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(CARD)
    shape.line.color.rgb = rgb(LINE)
    shape.line.width = Pt(1.0)

    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.18))
    strip.fill.solid()
    strip.fill.fore_color.rgb = rgb(accent)
    strip.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.22), width - Inches(0.32), Inches(0.34))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_style(run, title_size, color=INK, bold=True)

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.58), width - Inches(0.32), height - Inches(0.72))
    body_tf = body_box.text_frame
    body_tf.clear()
    body_tf.word_wrap = True
    for idx, line in enumerate(body):
        p2 = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p2.space_after = Pt(5)
        run2 = p2.add_run()
        run2.text = line
        set_run_style(run2, body_size, color=TEXT, name=FONT_ALT)
    return shape


def add_pill(slide, left: float, top: float, text: str, fill: str, *, width: float = Inches(1.12)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run_style(run, 10, color=INK, bold=True, name=FONT_ALT)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_chevron(slide, left: float, top: float, width: float = Inches(0.44), height: float = Inches(0.42), color: str = GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_flow_box(slide, left: float, top: float, width: float, height: float, title: str, subtitle: str | None = None, *, fill: str = CARD, line: str = LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_run_style(run, 16, color=INK, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        set_run_style(run2, 10, color=MUTED, name=FONT_ALT)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_route_row(slide, left: float, top: float, labels: list[tuple[str, str]], *, box_w: float = Inches(2.4)):
    current_left = left
    for idx, (title, accent) in enumerate(labels):
        add_flow_box(slide, current_left, top, box_w, Inches(0.8), title, None, fill=CARD, line=accent)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_left, top, box_w, Inches(0.08))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(accent)
        band.line.fill.background()
        current_left += box_w
        if idx != len(labels) - 1:
            add_chevron(slide, current_left + Inches(0.04), top + Inches(0.18), color=accent)
            current_left += Inches(0.56)


def add_table(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[str]], *, col_widths: list[float] | None = None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for idx, col_width in enumerate(col_widths):
            table.columns[idx].width = col_width

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(TEAL_SOFT)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_style(p.runs[0], 10, color=INK, bold=True)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(CARD if row_idx % 2 else "FBF8F2")
            cell.text = value
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(5)
            tf.margin_right = Pt(5)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_run_style(run, 9, color=TEXT, name=FONT_ALT)


def cover_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.55), Inches(1.3), Inches(0.1))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = rgb(TEAL)
    ribbon.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(7.4), Inches(1.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "HarborGate 到 HarborBeacon"
    set_run_style(run, 30, color=INK, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    run2.text = "北向入口、业务核心与南向域能力"
    set_run_style(run2, 17, color=TEAL, bold=True)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(10)
    run3 = p3.add_run()
    run3.text = "内部技术分享版 | 2026-04-24"
    set_run_style(run3, 11, color=MUTED, name=FONT_ALT)

    add_textbox(
        slide,
        Inches(0.7),
        Inches(3.0),
        Inches(5.8),
        Inches(1.6),
        [
            "这版 PPT 只做一件事：把一条请求从 IM 入口走到南向执行，再回到 IM 的全过程讲清楚。",
            "重点放在职责边界、冻结接口、北向入口面、南向域能力和三条典型链路。",
        ],
        size=16,
        color=TEXT,
    )

    add_card(
        slide,
        Inches(8.25),
        Inches(0.95),
        Inches(4.15),
        Inches(1.45),
        "HarborGate",
        ["IM adapters", "route_key / session_id", "platform delivery"],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(8.25),
        Inches(2.7),
        Inches(4.15),
        Inches(1.45),
        "HarborBeacon Core",
        ["task / session", "approval / artifact / audit", "planner / router / policy"],
        accent=TEAL,
        body_size=12,
    )
    add_card(
        slide,
        Inches(8.25),
        Inches(4.45),
        Inches(4.15),
        Inches(1.45),
        "Southbound Domains",
        ["HarborOS System Domain", "Home Device Domain", "Model Center + Retrieval"],
        accent=GOLD,
        body_size=12,
    )

    foot = slide.shapes.add_textbox(Inches(11.55), Inches(7.0), Inches(1.0), Inches(0.22))
    p = foot.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "01/15"
    set_run_style(run, 9, color=MUTED, name=FONT_ALT)


def slide_two(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "今天这 15 分钟讲什么", "先把地图看清，再讲接口、域能力和典型链路。", page=2)

    add_card(
        slide,
        Inches(0.68),
        Inches(1.45),
        Inches(5.45),
        Inches(3.65),
        "这次分享的 4 个问题",
        [
            "1. 为什么要拆成 HarborGate 和 HarborBeacon",
            "2. 一条请求在两个仓之间怎么流转",
            "3. HarborBeacon 的北向入口面和核心职责落在哪些模块",
            "4. 南向系统域、设备域和 model center 各自承接什么能力",
        ],
        accent=TEAL,
    )

    add_card(
        slide,
        Inches(6.45),
        Inches(1.45),
        Inches(2.75),
        Inches(1.75),
        "北向",
        [
            "面向 IM 平台、管理后台、页面和上层产品面的入口能力。",
        ],
        accent=RUST,
    )
    add_card(
        slide,
        Inches(9.45),
        Inches(1.45),
        Inches(2.85),
        Inches(1.75),
        "南向",
        [
            "面向 HarborOS、设备、协议适配器、本地模型和执行面的能力调用。",
        ],
        accent=GOLD,
    )

    add_card(
        slide,
        Inches(6.45),
        Inches(3.5),
        Inches(5.85),
        Inches(1.55),
        "怎么读后面的页面",
        [
            "上半场先看 HarborGate、HarborBeacon 和冻结边界；下半场看两条南向域和三条典型链路。",
            "每一页都按“方便投屏讲”的节奏压缩信息密度。",
        ],
        accent=NAVY,
        body_size=12,
    )

    add_pill(slide, Inches(0.75), Inches(5.55), "[已落地]", GREEN_SOFT)
    add_pill(slide, Inches(1.98), Inches(5.55), "[在途]", GOLD_SOFT)
    add_pill(slide, Inches(3.07), Inches(5.55), "[愿景]", RUST_SOFT)
    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.92),
        Inches(5.9),
        Inches(0.7),
        [
            "这套标签只用来区分现状、在途项和长期方向，讲的时候口径更稳。",
        ],
        size=11,
        color=MUTED,
    )


def slide_three(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "一条请求怎么从 HarborGate 走到 HarborBeacon", "先看全链路，再拆开每一段。", page=3)

    labels = [
        ("用户 / IM 平台", RUST),
        ("HarborGate", RUST),
        ("POST /api/tasks", TEAL),
        ("HarborBeacon 核心", TEAL),
        ("南向执行域", GOLD),
        ("POST /api/notifications/deliveries", TEAL),
        ("IM 回包 / 通知", RUST),
    ]
    left = Inches(0.72)
    top = Inches(2.2)
    widths = [Inches(1.4), Inches(1.4), Inches(1.6), Inches(1.8), Inches(1.55), Inches(2.05), Inches(1.55)]
    current = left
    for idx, ((label, accent), width) in enumerate(zip(labels, widths)):
        add_flow_box(slide, current, top, width, Inches(1.0), label, None, fill=CARD, line=accent)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, current, top, width, Inches(0.08))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(accent)
        band.line.fill.background()
        current += width
        if idx != len(labels) - 1:
            add_chevron(slide, current + Inches(0.05), top + Inches(0.26), color=accent)
            current += Inches(0.58)

    add_card(
        slide,
        Inches(0.85),
        Inches(4.2),
        Inches(3.6),
        Inches(1.45),
        "你可以抓住的 3 个锚点",
        [
            "HarborGate 负责平台接入、route key 和消息投递。",
            "HarborBeacon 负责任务状态、审批、artifact、audit 和南向编排。",
            "出站通知继续走 HarborGate，平台凭据也留在 HarborGate。",
        ],
        accent=TEAL,
        body_size=12,
    )
    add_card(
        slide,
        Inches(4.85),
        Inches(4.2),
        Inches(3.45),
        Inches(1.45),
        "三条固定 seam",
        [
            "`POST /api/tasks`",
            "`POST /api/notifications/deliveries`",
            "`GET /api/gateway/status`",
        ],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(8.65),
        Inches(4.2),
        Inches(3.55),
        Inches(1.45),
        "讲的时候顺手点一下",
        [
            "`X-Contract-Version: 1.5`",
            "service auth + RFC 3339 UTC",
            "`route_key` opaque, `resume_token` 跟着业务流走",
        ],
        accent=GOLD,
        body_size=12,
    )


def slide_four(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "两个仓各自负责什么", "这页讲清楚职责分工，后面每一页都更容易听。", page=4)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(4.55),
        Inches(4.65),
        "HarborGate",
        [
            "IM adapters、webhook、websocket / long connection、long-poll",
            "route_key / session_id 生命周期",
            "平台凭据保存、验证和 redacted status",
            "outbound delivery、payload 格式化、delivery retry",
            "setup portal、mobile onboarding、平台状态页",
        ],
        accent=RUST,
    )
    add_card(
        slide,
        Inches(8.0),
        Inches(1.45),
        Inches(4.55),
        Inches(4.65),
        "HarborBeacon",
        [
            "task / session lifecycle",
            "planner / router / policy / executor",
            "approval / artifact / audit / event",
            "HarborOS System Domain、Home Device Domain",
            "notification intent、model center、retrieval",
        ],
        accent=TEAL,
    )

    center = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.65), Inches(1.78), Inches(1.65), Inches(4.0))
    center.fill.solid()
    center.fill.fore_color.rgb = rgb("FBF6EF")
    center.line.color.rgb = rgb(LINE)
    center.line.width = Pt(1.0)
    center.text_frame.clear()
    center.text_frame.word_wrap = True
    items = [
        "冻结接口",
        "",
        "POST /api/tasks",
        "",
        "POST /api/notifications/deliveries",
        "",
        "GET /api/gateway/status",
    ]
    for idx, text in enumerate(items):
        p = center.text_frame.paragraphs[0] if idx == 0 else center.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_run_style(run, 12 if idx else 14, color=INK, bold=idx == 0)
    center.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(
        slide,
        Inches(0.9),
        Inches(6.2),
        Inches(11.4),
        Inches(0.42),
        [
            "这套拆分以后，IM 侧可以独立迭代平台接入，HarborBeacon 侧可以集中打磨任务、状态和南向能力；联调只盯住这 3 条 seam。",
        ],
        size=12,
        color=MUTED,
    )


def slide_five(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "冻结边界与契约细节", "把字段和约束讲清楚，团队协作会轻松很多。", page=5)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.3),
        Inches(3.7),
        Inches(1.0),
        "HarborGate -> HarborBeacon",
        ["`POST /api/tasks`", "标准任务入口，承接 inbound turn。"],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.3),
        Inches(3.7),
        Inches(1.0),
        "HarborBeacon -> HarborGate",
        ["`POST /api/notifications/deliveries`", "通知意图交给 HarborGate 执行投递。"],
        accent=TEAL,
        body_size=12,
    )
    add_card(
        slide,
        Inches(8.85),
        Inches(1.3),
        Inches(3.55),
        Inches(1.0),
        "HarborBeacon <- HarborGate",
        ["`GET /api/gateway/status`", "管理面消费脱敏平台状态。"],
        accent=GOLD,
        body_size=12,
    )

    add_card(
        slide,
        Inches(0.75),
        Inches(2.65),
        Inches(3.78),
        Inches(2.45),
        "TaskRequest",
        [
            "`task_id / trace_id / step_id`",
            "`source.channel / surface / conversation_id / user_id / session_id / route_key`",
            "`intent.domain / action / raw_text`",
            "`args / autonomy.level / message.message_id / chat_type / mentions / attachments`",
        ],
        accent=TEAL,
        body_size=11,
    )
    add_card(
        slide,
        Inches(4.78),
        Inches(2.65),
        Inches(3.78),
        Inches(2.45),
        "TaskResponse",
        [
            "`status / executor_used / risk_level`",
            "`result.message / data / artifacts / events / next_actions`",
            "`audit_ref / missing_fields / prompt / resume_token`",
            "`completed / needs_input / failed` 三种状态贯穿多轮流程",
        ],
        accent=RUST,
        body_size=11,
    )
    add_card(
        slide,
        Inches(8.81),
        Inches(2.65),
        Inches(3.59),
        Inches(2.45),
        "Notification + GatewayStatus",
        [
            "`notification_id / trace_id / destination.route_key / delivery.mode / idempotency_key`",
            "`GatewayStatusResponse` 返回平台启用状态、连接状态、能力和 redacted status",
            "`route_key` 只当路由句柄来用",
        ],
        accent=GOLD,
        body_size=11,
    )

    add_card(
        slide,
        Inches(0.75),
        Inches(5.35),
        Inches(11.65),
        Inches(1.15),
        "讲这页时顺手补的 6 个口径",
        [
            "`X-Contract-Version: 1.5` 进入所有冻结接口；service-to-service auth 单独校验；跨仓时间戳用 RFC 3339 UTC；`task_id` 支持 replay / conflict；request rejection 走非 200 shared error envelope；accepted delivery failure 用 `HTTP 200 + ok=false`。",
        ],
        accent=NAVY,
        body_size=12,
    )


def slide_six(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "HarborGate 北向传输层", "平台入口、route 语义和运维入口都集中在这一侧。", page=6)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(3.55),
        Inches(1.75),
        "当前平台覆盖",
        [
            "live: `feishu / weixin / webhook`",
            "placeholder: `telegram / discord / slack / whatsapp / signal / email / wecom`",
        ],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(3.45),
        Inches(3.55),
        Inches(1.7),
        "接入模式",
        [
            "webhook",
            "websocket / long connection",
            "long-poll",
        ],
        accent=GOLD,
        body_size=12,
    )
    add_card(
        slide,
        Inches(4.55),
        Inches(1.45),
        Inches(3.75),
        Inches(3.7),
        "运行时语义",
        [
            "`route_key`：平台路由句柄",
            "`session_id`：会话级别标识",
            "`resume_token`：按 chat 保存，后续 turn 带回 HarborBeacon",
            "`task_id / trace_id`：从 inbound event identity 派生稳定值",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(8.55),
        Inches(1.45),
        Inches(3.8),
        Inches(3.7),
        "运维入口",
        [
            "`/setup /setup/qr /setup/qr.svg`",
            "`/api/setup/status`",
            "`/api/setup/feishu/configure`",
            "`/api/gateway/status`",
            "`/api/notifications/deliveries`",
        ],
        accent=NAVY,
    )

    add_textbox(
        slide,
        Inches(0.82),
        Inches(5.55),
        Inches(11.2),
        Inches(0.55),
        [
            "讲到 `GET /api/gateway/status` 时可以补一句：HarborGate 会把平台状态收敛成 `not_configured / configured_placeholder / live` 这一类统一口径，管理面读的是脱敏状态和能力摘要。",
        ],
        size=11,
        color=MUTED,
    )


def slide_seven(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "HarborBeacon 北向入口面", "多个入口共用同一套 task / runtime / approval / audit 底座。", page=7)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.5),
        Inches(3.25),
        Inches(2.0),
        "assistant-task-api",
        [
            "`POST /api/tasks`",
            "`GET /healthz`",
            "HarborGate 进来的标准 turn 会先落到这里。",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(4.2),
        Inches(1.5),
        Inches(4.1),
        Inches(2.0),
        "agent-hub-admin-api",
        [
            "`/api/state /api/account-management`",
            "`/api/tasks/approvals /api/discovery/scan /api/devices/manual`",
            "`/api/cameras/* /api/share-links /api/admin/notification-targets`",
            "`/api/models/endpoints /api/models/policies /admin/models`",
        ],
        accent=RUST,
        body_size=11,
    )
    add_card(
        slide,
        Inches(8.55),
        Inches(1.5),
        Inches(3.8),
        Inches(2.0),
        "harbor-model-api",
        [
            "`GET /healthz`",
            "`POST /v1/chat/completions`",
            "`POST /v1/embeddings`",
            "模型能力作为服务面挂在 HarborBeacon 侧。",
        ],
        accent=GOLD,
        body_size=11,
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(3.9),
        Inches(11.6),
        Inches(1.45),
        "Harbor Assistant / WebUI / live pages",
        [
            "Harbor Assistant、绑定页、live view、share page 都属于产品面；它们负责展示、配置和交互，任务状态和审批链继续回到同一套 runtime 里。",
        ],
        accent=NAVY,
        body_size=12,
    )
    add_route_row(
        slide,
        Inches(1.25),
        Inches(5.8),
        [
            ("入口很多", NAVY),
            ("状态一套", TEAL),
            ("审计一套", TEAL),
            ("通知一套", RUST),
        ],
        box_w=Inches(2.25),
    )


def slide_eight(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "HarborBeacon 业务核心", "这一层承接任务、会话、审批、artifact、audit 和通知意图。", page=8)

    steps = [
        "入口接收",
        "幂等判断",
        "会话恢复",
        "planner / router / policy",
        "执行或触发审批",
        "写入 task run / artifact / event",
        "返回 TaskResponse / 生成通知意图",
    ]
    top = Inches(1.55)
    current = Inches(0.8)
    for idx, step in enumerate(steps):
        add_flow_box(slide, current, top, Inches(1.55), Inches(0.85), step, None, fill=CARD, line=TEAL if idx < 4 else RUST)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, current, top, Inches(1.55), Inches(0.08))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(TEAL if idx < 4 else RUST)
        band.line.fill.background()
        current += Inches(1.55)
        if idx != len(steps) - 1:
            add_chevron(slide, current + Inches(0.02), top + Inches(0.21), width=Inches(0.34), color=GOLD)
            current += Inches(0.42)

    add_card(
        slide,
        Inches(0.82),
        Inches(3.15),
        Inches(5.35),
        Inches(2.25),
        "当前这套核心已经持有的业务真相",
        [
            "business session state",
            "resumable workflow state",
            "approvals / artifacts / audit trail",
            "business conversation continuity",
            "notification intent、media sessions、share links",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(6.45),
        Inches(3.15),
        Inches(5.87),
        Inches(2.25),
        "你可以顺手点一下的持久化对象",
        [
            "`conversations / sessions / task_runs / task_steps`",
            "`artifacts / approvals / events`",
            "`media_sessions / share_links`",
            "高风险动作会变成 `needs_input` 或 pending approval，再由后续 turn 继续完成。",
        ],
        accent=RUST,
        body_size=12,
    )


def slide_nine(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "南向 A：HarborOS System Domain", "固定优先级已经收口到控制面优先。", page=9)

    add_route_row(
        slide,
        Inches(0.95),
        Inches(1.45),
        [
            ("Middleware API", TEAL),
            ("MidCLI", RUST),
            ("Browser / MCP fallback", GOLD),
        ],
        box_w=Inches(3.35),
    )

    add_card(
        slide,
        Inches(0.82),
        Inches(2.55),
        Inches(5.7),
        Inches(2.55),
        "已落地系统能力",
        [
            "`service.query`",
            "`service.control`",
            "`files.copy`",
            "`files.move`",
            "`files.list`",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(6.75),
        Inches(2.55),
        Inches(5.55),
        Inches(2.55),
        "讲这页时顺手带上的边界",
        [
            "高风险系统动作会走 approval gate。",
            "`Browser / MCP` 只放在 fallback 位，不接管 HarborOS ownership。",
            "`files.stat / files.read_text` 留在 helper scope，方便 framework 使用，不扩成 HarborOS 产品面。",
        ],
        accent=RUST,
        body_size=12,
    )

    add_textbox(
        slide,
        Inches(0.85),
        Inches(5.55),
        Inches(11.0),
        Inches(0.5),
        [
            "这一页的关键句：系统域先走 HarborOS 自己的接口面，回退路径是为了可用性和兼容性。",
        ],
        size=11,
        color=MUTED,
    )


def slide_ten(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "南向 B：Home Device Domain", "设备域单独成线，摄像头协议和媒体链路都在这里收口。", page=10)

    add_route_row(
        slide,
        Inches(0.78),
        Inches(1.45),
        [
            ("Native Adapter", TEAL),
            ("LAN Bridge", RUST),
            ("HarborOS Connector", GOLD),
            ("Cloud / MCP", NAVY),
        ],
        box_w=Inches(2.5),
    )

    add_card(
        slide,
        Inches(0.8),
        Inches(2.65),
        Inches(3.95),
        Inches(2.65),
        "协议与适配",
        [
            "`ONVIF`",
            "`SSDP`",
            "`mDNS`",
            "`RTSP probe`",
            "vendor-cloud bridge",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(4.95),
        Inches(2.65),
        Inches(3.55),
        Inches(2.65),
        "domain actions",
        [
            "`discover / list / get / update`",
            "`snapshot / open_stream / ptz`",
            "`camera.scan / connect / snapshot / share_link / analyze`",
        ],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(8.72),
        Inches(2.65),
        Inches(3.55),
        Inches(2.65),
        "这一域已经带出来的平台问题",
        [
            "设备发现和接入",
            "媒体 / 控制分离",
            "live view / share link / capture / analyze",
            "人机补参和长任务状态",
        ],
        accent=GOLD,
        body_size=12,
    )


def slide_eleven(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "能力补能：Model Center + Multimodal Retrieval", "这条线贴着 HarborBeacon 核心走，方便统一治理和路由。", page=11)

    add_card(
        slide,
        Inches(0.78),
        Inches(1.45),
        Inches(5.4),
        Inches(2.65),
        "当前范围",
        [
            "`document + image + OCR` 已经接到同一条检索主线。",
            "citation 和 reply pack 继续由 HarborBeacon 打包。",
            "`harbor-model-api` 提供 `/healthz`、chat completions 和 embeddings。",
            "管理面已经有 endpoints / policies 的读写和测试入口。",
        ],
        accent=TEAL,
        body_size=12,
    )
    add_card(
        slide,
        Inches(6.45),
        Inches(1.45),
        Inches(5.87),
        Inches(2.65),
        "4 个策略位",
        [
            "`retrieval.ocr`：优先本地 tesseract",
            "`retrieval.embed`：优先本地 OpenAI-compatible endpoint",
            "`retrieval.answer`：local-first，可云端补能",
            "`retrieval.vision_summary`：已占位，等 VLM 配置后提速",
        ],
        accent=RUST,
        body_size=12,
    )
    add_card(
        slide,
        Inches(0.78),
        Inches(4.45),
        Inches(11.55),
        Inches(1.5),
        "讲这页时可以这样收口",
        [
            "model center 负责把 OCR、embedding、answer 和 vision summary 变成可治理、可切换、可观测的能力槽位；更完整的 audio / video / full multimodal 还在后续阶段。",
        ],
        accent=NAVY,
        body_size=12,
    )


def slide_twelve(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "三条典型链路", "用故事把前面的结构串起来，现场讲解会更顺。", page=12)

    add_card(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(3.8),
        Inches(4.85),
        "案例 1：IM 查询系统状态",
        [
            "1. 用户在飞书或微信发出状态查询。",
            "2. HarborGate 归一化消息，带上 `task_id / trace_id / route_key`。",
            "3. `POST /api/tasks` 进入 HarborBeacon。",
            "4. HarborBeacon 路由到 system domain。",
            "5. `Middleware API -> MidCLI -> fallback` 选路执行。",
            "6. 结果回到 HarborGate，再映射成平台 reply。",
        ],
        accent=TEAL,
        body_size=11,
    )
    add_card(
        slide,
        Inches(4.77),
        Inches(1.55),
        Inches(3.8),
        Inches(4.85),
        "案例 2：高风险操作的审批恢复",
        [
            "1. 用户发起 restart、move 之类的高风险动作。",
            "2. HarborBeacon policy 判断需要 approval。",
            "3. 返回 `status=needs_input`、`prompt`、`resume_token` 或审批 ticket。",
            "4. HarborGate 把提示送回用户，管理面也能看到待审批任务。",
            "5. 审批通过后，HarborBeacon 按 approval context replay / resume。",
            "6. 最终结果再走同步 reply 或后续 notification。",
        ],
        accent=RUST,
        body_size=11,
    )
    add_card(
        slide,
        Inches(8.82),
        Inches(1.55),
        Inches(3.8),
        Inches(4.85),
        "案例 3：摄像头从发现到回送",
        [
            "1. 用户或管理面发起扫描。",
            "2. HarborBeacon 走 `ONVIF / SSDP / mDNS / RTSP probe`。",
            "3. 返回候选设备列表，必要时进入补参。",
            "4. 接入后写入 registry。",
            "5. 抓拍、分析、分享继续走 device / media / vision 链。",
            "6. 结果变成 artifact、share link、event，再经 HarborGate 送回用户。",
        ],
        accent=GOLD,
        body_size=11,
    )


def slide_thirteen(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Home Agent Hub 为什么值得单独讲", "它把平台抽象推进到了可跑、可审计、可回送的程度。", page=13)

    add_card(
        slide,
        Inches(0.78),
        Inches(1.5),
        Inches(5.15),
        Inches(4.4),
        "这个垂直域已经验证出来的平台能力",
        [
            "artifact envelope：文本、图片、录像、预览链接、分享链接、动作卡片",
            "long-running task：扫描、抓流、录像、视觉分析都要有状态管理",
            "entity resolution：设备、房间、对象、绑定用户要能从自然语言里抽出来",
            "human-in-the-loop：选号、补密码、继续执行要走标准机制",
            "sidecar worker pattern：ffmpeg、YOLO、本地播放器、后续工具 worker 都要接入治理面",
        ],
        accent=TEAL,
        body_size=12,
    )

    labels = [
        ("Artifact", TEAL_SOFT, Inches(6.55), Inches(1.9)),
        ("Long Task", RUST_SOFT, Inches(8.55), Inches(3.05)),
        ("HITL", GOLD_SOFT, Inches(10.4), Inches(1.9)),
        ("Entity", NAVY_SOFT, Inches(7.45), Inches(4.3)),
        ("Workers", GREEN_SOFT, Inches(9.6), Inches(4.3)),
    ]
    for text, fill, left, top in labels:
        add_pill(slide, left, top, text, fill, width=Inches(1.55))

    add_card(
        slide,
        Inches(6.2),
        Inches(5.85),
        Inches(5.9),
        Inches(0.95),
        "一句话收口",
        [
            "Home Agent Hub 让 HarborBeacon 的平台抽象从文档定义走到了真实域能力：有任务、有状态、有审计，也有完整回送链路。",
        ],
        accent=RUST,
        body_size=12,
    )


def slide_fourteen(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "附录 A：北向接口 / 入口矩阵", "会后留档时，这张表最方便快速翻。", page=14)

    headers = ["层级", "入口 / 接口", "主要路径", "角色定位", "状态"]
    rows = [
        ["HarborGate", "平台消息入口", "POST /messages/<platform>", "adapter 归一化入口", "[已落地]"],
        ["HarborGate", "setup portal", "/setup /setup/qr /api/setup/status", "北向配置与启动引导", "[已落地]"],
        ["HarborGate", "gateway status", "GET /api/gateway/status", "脱敏状态和能力摘要", "[已落地]"],
        ["HarborGate -> HarborBeacon", "task ingress seam", "POST /api/tasks", "冻结跨仓任务入口", "[已落地]"],
        ["HarborBeacon", "assistant-task-api", "POST /api/tasks, GET /healthz", "统一任务入口", "[已落地]"],
        ["HarborBeacon", "agent-hub-admin-api", "/api/state /api/tasks/approvals /api/cameras/* /api/models/*", "管理面与产品面入口", "[已落地]"],
        ["HarborBeacon", "harbor-model-api", "/healthz /v1/chat/completions /v1/embeddings", "模型能力服务面", "[已落地]"],
        ["HarborBeacon -> HarborGate", "notification seam", "POST /api/notifications/deliveries", "冻结跨仓通知投递意图", "[已落地]"],
    ]
    add_table(
        slide,
        Inches(0.62),
        Inches(1.55),
        Inches(12.05),
        Inches(5.7),
        headers,
        rows,
        col_widths=[Inches(1.45), Inches(2.0), Inches(3.05), Inches(4.25), Inches(1.1)],
    )


def slide_fifteen(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "附录 B：南向能力 / 路由矩阵", "这张表适合回答“这个能力到底落在哪条域线里”。", page=15)

    headers = ["域", "能力组", "具体动作 / 能力", "首选路由", "回退 / 补位", "状态"]
    rows = [
        ["HarborOS System", "服务查询", "service.query", "Middleware API", "MidCLI -> Browser/MCP fallback", "[已落地]"],
        ["HarborOS System", "服务控制", "service.control", "Middleware API", "MidCLI -> Browser/MCP fallback", "[已落地]"],
        ["HarborOS System", "文件能力", "files.list / copy / move", "Middleware API", "MidCLI -> constrained fallback", "[已落地]"],
        ["HarborOS System", "helper scope", "files.stat / files.read_text", "helper only", "保持 preview 范围", "[已落地]"],
        ["Home Device", "发现与接入", "discover / list / get / update / camera.scan / connect", "Native Adapter", "LAN Bridge -> HarborOS Connector -> Cloud/MCP", "[已落地]"],
        ["Home Device", "媒体与控制", "snapshot / open_stream / ptz / share_link / analyze", "RTSP / native control", "vendor-cloud bridge / sidecar", "[已落地]"],
        ["Model Center", "OCR / embed / answer", "retrieval.ocr / embed / answer", "local-first", "cloud fallback", "[已落地]"],
        ["Model Center", "vision summary", "retrieval.vision_summary", "configured VLM", "degraded until ready", "[在途]"],
    ]
    add_table(
        slide,
        Inches(0.55),
        Inches(1.55),
        Inches(12.2),
        Inches(5.7),
        headers,
        rows,
        col_widths=[Inches(1.5), Inches(1.6), Inches(2.6), Inches(1.55), Inches(3.8), Inches(0.95)],
    )


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    slide_five(prs)
    slide_six(prs)
    slide_seven(prs)
    slide_eight(prs)
    slide_nine(prs)
    slide_ten(prs)
    slide_eleven(prs)
    slide_twelve(prs)
    slide_thirteen(prs)
    slide_fourteen(prs)
    slide_fifteen(prs)
    return prs


def main():
    prs = build_deck()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
